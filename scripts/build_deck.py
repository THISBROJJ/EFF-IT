"""Build EFF-IT presentation as a .pptx file.

Usage:
    python scripts/build_deck.py [output_path]

Default output: docs/EFF-IT.pptx
"""

from __future__ import annotations

import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt


REPO_ROOT = Path(__file__).resolve().parent.parent
SCREENSHOTS_DIR = REPO_ROOT / "pptx_refs"
ASSETS_DIR = REPO_ROOT / "docs" / "assets" / "case-study"


def find_screenshot(time_prefix: str, suffix: str = "") -> Path:
    """Locate a screenshot by its HH.MM.SS prefix.

    Filenames contain a U+202F (narrow no-break space) before 'AM' which
    makes direct lookup brittle — globbing is the reliable path.
    """
    pattern = f"Screenshot 2026-05-22 at {time_prefix}*{suffix}*.png"
    matches = sorted(SCREENSHOTS_DIR.glob(pattern))
    if not matches:
        raise FileNotFoundError(f"No screenshot for {time_prefix!r}")
    return matches[0]


def prep_image(time_prefix: str, slug: str, max_width: int = 1600) -> Path:
    """Copy + downscale a source screenshot into docs/assets/case-study/.

    Returns the path to the prepared JPEG.
    """
    src = find_screenshot(time_prefix)
    ASSETS_DIR.mkdir(parents=True, exist_ok=True)
    dst = ASSETS_DIR / f"{slug}.jpg"
    if dst.exists() and dst.stat().st_mtime >= src.stat().st_mtime:
        return dst
    img = Image.open(src).convert("RGB")
    if img.width > max_width:
        ratio = max_width / img.width
        img = img.resize((max_width, int(img.height * ratio)), Image.LANCZOS)
    img.save(dst, "JPEG", quality=85, optimize=True)
    return dst


# Palette
BG = RGBColor(0x0F, 0x14, 0x1A)
PANEL = RGBColor(0x18, 0x1F, 0x2A)
ACCENT = RGBColor(0x6E, 0xE7, 0xB7)
ACCENT_DIM = RGBColor(0x34, 0xA8, 0x7E)
TEXT = RGBColor(0xE8, 0xEC, 0xF1)
MUTED = RGBColor(0x9A, 0xA3, 0xAF)
DANGER = RGBColor(0xF8, 0x72, 0x71)
WARN = RGBColor(0xF5, 0xC2, 0x49)


def set_bg(slide, prs):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg


def add_textbox(slide, left, top, width, height, text, *,
                size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT,
                font="Segoe UI", mono=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = "Consolas" if mono else font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return box


def add_rich_paragraphs(slide, left, top, width, height, paragraphs,
                        font="Segoe UI"):
    """paragraphs: list of list of (text, dict-of-run-formatting)."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, runs in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        for text, fmt in runs:
            r = p.add_run()
            r.text = text
            r.font.name = fmt.get("font", font)
            r.font.size = Pt(fmt.get("size", 18))
            r.font.bold = fmt.get("bold", False)
            r.font.italic = fmt.get("italic", False)
            r.font.color.rgb = fmt.get("color", TEXT)
    return box


def add_panel(slide, left, top, width, height, color=PANEL):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = ACCENT_DIM
    shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    if shape.text_frame.paragraphs:
        shape.text_frame.paragraphs[0].text = ""
    return shape


def add_accent_bar(slide, left, top, width=Inches(0.08), height=Inches(0.5)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    bar.shadow.inherit = False
    return bar


def add_slide_header(slide, title, subtitle=None):
    add_accent_bar(slide, Inches(0.5), Inches(0.45), height=Inches(0.55))
    add_textbox(slide, Inches(0.75), Inches(0.35), Inches(12), Inches(0.7),
                title, size=32, bold=True, color=TEXT)
    if subtitle:
        add_textbox(slide, Inches(0.75), Inches(1.0), Inches(12), Inches(0.4),
                    subtitle, size=14, color=MUTED)


def add_footer(slide, idx, total):
    add_textbox(slide, Inches(0.5), Inches(7.0), Inches(6), Inches(0.3),
                "EFF-IT  —  A Claude Code harness", size=10, color=MUTED)
    add_textbox(slide, Inches(11.5), Inches(7.0), Inches(1.8), Inches(0.3),
                f"{idx} / {total}", size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def add_notes(slide, notes):
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = notes


def add_pill(slide, left, top, height, label, *, color=ACCENT, fill=PANEL,
             padding=0.25, size=14, bold=True):
    """A rounded chip with auto-sized width based on label length."""
    char_w = size * 0.075  # rough em width estimate in inches
    width = Inches(char_w * len(label) + padding * 2)
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    left, top, width, height)
    pill.fill.solid()
    pill.fill.fore_color.rgb = fill
    pill.line.color.rgb = color
    pill.line.width = Pt(1.25)
    pill.shadow.inherit = False
    tf = pill.text_frame
    tf.margin_left = Inches(padding)
    tf.margin_right = Inches(padding)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.name = "Segoe UI"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return pill, width


def add_card(slide, left, top, width, height, *,
             accent_color=ACCENT, title=None, body_lines=None,
             title_size=16, body_size=13):
    """Reusable card: panel with a colored top bar + title + body lines."""
    add_panel(slide, left, top, width, height)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()
    bar.shadow.inherit = False
    if title:
        add_textbox(slide, left + Inches(0.2), top + Inches(0.18),
                    width - Inches(0.4), Inches(0.4),
                    title, size=title_size, bold=True, color=accent_color)
    if body_lines:
        y_off = 0.65
        for line in body_lines:
            add_textbox(slide, left + Inches(0.25), top + Inches(y_off),
                        width - Inches(0.5), Inches(0.45),
                        f"•  {line}", size=body_size, color=TEXT)
            y_off += 0.4


def add_table(slide, left, top, width, height, rows, *,
              header=True, col_widths=None, font_size=14):
    nrows = len(rows)
    ncols = len(rows[0])
    table_shape = slide.shapes.add_table(nrows, ncols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        total = sum(col_widths)
        for i, frac in enumerate(col_widths):
            table.columns[i].width = int(width * frac / total)
    for r_idx, row in enumerate(rows):
        for c_idx, cell_text in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.fill.solid()
            if header and r_idx == 0:
                cell.fill.fore_color.rgb = ACCENT_DIM
            else:
                cell.fill.fore_color.rgb = PANEL if r_idx % 2 else BG
            tf = cell.text_frame
            tf.margin_left = Inches(0.1)
            tf.margin_right = Inches(0.1)
            tf.margin_top = Inches(0.05)
            tf.margin_bottom = Inches(0.05)
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = cell_text
            run.font.name = "Segoe UI"
            run.font.size = Pt(font_size)
            run.font.bold = (header and r_idx == 0)
            run.font.color.rgb = TEXT if not (header and r_idx == 0) else BG
    return table


# ---------- slide builders ----------

def slide_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, prs)
    # Big EFF-IT
    add_textbox(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(1.5),
                "EFF-IT", size=96, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.5), Inches(3.5), Inches(12.3), Inches(0.6),
                "A Claude Code harness for vibe-driven software delivery",
                size=24, color=TEXT, align=PP_ALIGN.CENTER)
    # Tagline panel
    add_panel(s, Inches(2.5), Inches(4.6), Inches(8.3), Inches(1.0))
    add_textbox(s, Inches(2.5), Inches(4.75), Inches(8.3), Inches(0.7),
                "“You know what? EFF-IT. Let the agents handle it.”",
                size=20, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
                "lizhang  ·  2026", size=12, color=MUTED, align=PP_ALIGN.CENTER)
    add_notes(s, (
        "Open by saying the name out loud, then the tagline. Pause for the laugh.\n\n"
        "Set expectations in two sentences:\n"
        "\"This is a presentation about a scaffold I built on top of Claude Code. "
        "It takes the 'vibe coding' experience — describe what you want, let the AI build it — "
        "and wraps it with the structure you'd normally have to remember to add yourself: "
        "specs, audits, security review, checkpoints, the works.\"\n\n"
        "Don't dive into architecture yet. The next slide explains why this exists at all. "
        "Resist the urge to list features here — you have ten slides for that."
    ))
    return s


def slide_what_is_vibe(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, prs)
    add_slide_header(s, "What is \"vibe coding\"?")
    bullets = [
        "You describe intent in plain language",
        "The agent reads the repo, writes code, runs tests, opens the PR",
        "You stay in the driver's seat — but you're not typing every line",
    ]
    y = 1.8
    for b in bullets:
        # bullet dot
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.9), Inches(y + 0.18),
                                  Inches(0.15), Inches(0.15))
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT
        dot.line.fill.background()
        add_textbox(s, Inches(1.2), Inches(y), Inches(11.5), Inches(0.5),
                    b, size=20, color=TEXT)
        y += 0.7

    add_panel(s, Inches(0.75), Inches(4.4), Inches(11.8), Inches(0.9))
    add_rich_paragraphs(s, Inches(1.0), Inches(4.55), Inches(11.5), Inches(0.7), [
        [("The promise: ", {"bold": True, "color": ACCENT, "size": 18}),
         ("ship features at the speed of thought.", {"size": 18})],
    ])

    add_panel(s, Inches(0.75), Inches(5.45), Inches(11.8), Inches(1.0))
    add_rich_paragraphs(s, Inches(1.0), Inches(5.6), Inches(11.5), Inches(0.8), [
        [("The reality (without guardrails): ", {"bold": True, "color": DANGER, "size": 18}),
         ("drift, half-done work, broken tests, leaked secrets.", {"size": 18})],
    ])
    add_notes(s, (
        "Make sure the room agrees on what 'vibe coding' means before you critique it. "
        "The term is fresh enough that half the audience will have heard it and half won't. "
        "Define it cleanly: you describe intent, the agent executes. That's it.\n\n"
        "The pivot in this slide is the last line — 'the reality without guardrails.' "
        "Land that beat hard. This is the problem statement for the whole talk. "
        "If the audience nods here, the rest of the deck writes itself.\n\n"
        "Concrete examples to drop if energy is low:\n"
        "  - 'I've watched Claude edit a failing test to make it pass.'\n"
        "  - 'I've watched it commit a .env file with a real API key.'\n"
        "  - 'I've watched it forget what it was building three messages in.'\n"
        "Pick one. Don't list all three.\n\n"
        "Transition: 'So I built something to fix that.'"
    ))


def slide_problem(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, prs)
    add_slide_header(s, "The problem with raw vibe coding")
    rows = [
        ["What happens", "Why it's painful"],
        ["Agent forgets context mid-feature", "Start over, lose state"],
        ["Tests get \"adjusted\" to pass", "False green builds"],
        ["Scope creeps silently", "Reviewer surprise"],
        ["Secrets land in commits", "Rotate, redeploy, apologize"],
        ["No record of why anything happened", "Can't debug the agent itself"],
    ]
    add_table(s, Inches(0.75), Inches(1.7), Inches(11.8), Inches(3.5), rows,
              col_widths=[1, 1], font_size=16)
    add_panel(s, Inches(0.75), Inches(5.6), Inches(11.8), Inches(1.1))
    add_rich_paragraphs(s, Inches(1.0), Inches(5.75), Inches(11.5), Inches(0.9), [
        [("Raw Claude Code is a ", {"size": 20}),
         ("power tool", {"bold": True, "color": WARN, "size": 20}),
         (". EFF-IT is the ", {"size": 20}),
         ("jig", {"bold": True, "color": ACCENT, "size": 20}),
         (" that makes it safe.", {"size": 20})],
    ])
    add_notes(s, (
        "Walk the table top-to-bottom. One sentence per row, no more:\n\n"
        "  1. 'Sessions die. Without a checkpoint, you start from scratch.'\n"
        "  2. 'Agents are great at making tests pass — including by editing the test.'\n"
        "  3. 'Ask for a bug fix, get a 40-file refactor.'\n"
        "  4. 'Secrets in commits aren't theoretical. It's happened to me twice.'\n"
        "  5. 'When the agent does something weird, you want a transcript. "
        "Raw Claude Code doesn't keep one by default.'\n\n"
        "Land the closing line — the power-tool / jig metaphor — slowly. "
        "This is the mental model for the whole rest of the deck. "
        "The harness is not a replacement for Claude Code; it's the fixture "
        "that holds the work piece steady while the tool does its job.\n\n"
        "If anyone asks 'why not just be more careful?' — that's the exact "
        "question the next slide answers."
    ))


def slide_what_is_eff_it(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, prs)
    add_slide_header(s, "What EFF-IT is", "A drop-in scaffold for any repo")
    items = [
        (".claude/commands/", "slash-command workflows (/run, /fast-lane, /resume)"),
        (".claude/agents/", "12 single-purpose subagents"),
        (".claude/hooks/", "lifecycle scripts (logging, secrets scan, test immutability)"),
        ("sessions/<run_id>/", "self-contained per-run artifacts + checkpoint"),
        ("security/profiles/", "app-type threat models loaded per run"),
    ]
    y = 1.7
    for path, desc in items:
        add_panel(s, Inches(0.75), Inches(y), Inches(11.8), Inches(0.65))
        add_textbox(s, Inches(1.0), Inches(y + 0.12), Inches(3.8), Inches(0.5),
                    path, size=15, bold=True, color=ACCENT, mono=True)
        add_textbox(s, Inches(4.9), Inches(y + 0.12), Inches(7.8), Inches(0.5),
                    desc, size=15, color=TEXT)
        y += 0.8

    add_textbox(s, Inches(0.75), Inches(5.9), Inches(11.8), Inches(0.4),
                "Stack: Bash + Markdown + GitHub Actions.   No runtime.",
                size=16, bold=True, color=WARN)
    add_textbox(s, Inches(0.75), Inches(6.3), Inches(11.8), Inches(0.4),
                "It targets whatever stack your project already uses.",
                size=14, color=MUTED)
    add_notes(s, (
        "Now you can finally show the goods. Five bullets, each pointing at a directory in the repo. "
        "Don't read every bullet — point at the slide and call out the three that matter most:\n\n"
        "  - 'Commands are how YOU talk to it — slash commands you type.'\n"
        "  - 'Agents are how IT talks to itself — each one does one thing.'\n"
        "  - 'Hooks are the safety rails that fire whether the agent remembers to or not.'\n\n"
        "The last paragraph is important: this isn't a framework that owns your stack. "
        "It's pure Bash and Markdown. It doesn't care if your project is Python, Go, Rust, or TypeScript. "
        "That portability is the reason it's worth adopting.\n\n"
        "If someone asks 'is this open source / can I copy it?' — yes, that's exactly the point. "
        "The README has the copy-paste steps. (Slide 11 covers adoption.)"
    ))


def slide_entry_points(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, prs)
    add_slide_header(s, "Three entry points for three vibes")
    cards = [
        ("/run [idea]", "New feature, vague brief", "Full pipeline, idea → PR", ACCENT),
        ("/fast-lane [desc]", "You know what to build", "Skip the spec phase", WARN),
        ("/resume <run_id>", "Session crashed", "Pick up at last checkpoint", DANGER),
    ]
    x = 0.5
    card_w = 4.1
    for label, when, what, color in cards:
        add_panel(s, Inches(x), Inches(1.8), Inches(card_w), Inches(3.4))
        # accent bar
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.8),
                                   Inches(card_w), Inches(0.12))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        add_textbox(s, Inches(x + 0.2), Inches(2.1), Inches(card_w - 0.4), Inches(0.6),
                    label, size=20, bold=True, color=color, mono=True)
        add_textbox(s, Inches(x + 0.2), Inches(2.9), Inches(card_w - 0.4), Inches(0.4),
                    "WHEN", size=11, bold=True, color=MUTED)
        add_textbox(s, Inches(x + 0.2), Inches(3.2), Inches(card_w - 0.4), Inches(0.8),
                    when, size=16, color=TEXT)
        add_textbox(s, Inches(x + 0.2), Inches(4.0), Inches(card_w - 0.4), Inches(0.4),
                    "WHAT", size=11, bold=True, color=MUTED)
        add_textbox(s, Inches(x + 0.2), Inches(4.3), Inches(card_w - 0.4), Inches(0.8),
                    what, size=16, color=TEXT)
        x += card_w + 0.2

    add_panel(s, Inches(0.75), Inches(5.6), Inches(11.8), Inches(1.0))
    add_textbox(s, Inches(1.0), Inches(5.75), Inches(11.5), Inches(0.8),
                "All three land in the same pipeline. /resume exists because long runs will get interrupted — and starting over is the worst feeling in vibe coding.",
                size=14, color=TEXT)
    add_notes(s, (
        "Three commands, three moods. Walk the row:\n\n"
        "  - '/run is the front door. Pitch me an idea — even a bad one — and it'll "
        "interrogate it into a spec before writing a line of code.'\n"
        "  - '/fast-lane is for when you don't want the Socratic dialogue. "
        "You've thought about it, you know what you want, just go.'\n"
        "  - '/resume is the unsung hero. Long agent runs die. Network blips, "
        "context fills up, you close the laptop. /resume reads checkpoint.json and "
        "drops you back in at the exact stage where things stopped.'\n\n"
        "The last paragraph is the emotional beat — anyone who's done serious agent work "
        "has lost a run at minute 45 and felt that specific kind of despair.\n\n"
        "Transition: 'What actually happens inside /run? Next slide.'"
    ))


def slide_pipeline(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, prs)
    add_slide_header(s, "The pipeline (the whole show)", "/run from idea to PR")

    stages = [
        ("idea-interrogator", "Socratic Q&A until spec is concrete"),
        ("spec-drafter", "→ SPEC.md"),
        ("git-expert", "feat/<slug> branch"),
        ("orchestrator", "→ PLAN.md (task decomposition)"),
        ("concern-resolver", "→ SECURITY_CONCERNS.md"),
        ("architect", "→ ARCHITECTURE.md"),
        ("implementation-loop", "unit-test-writer → coder → test-runner  (≤ 5x)"),
        ("karen", "PASS / PARTIAL / FAIL audit vs SPEC"),
        ("evaluate-run", "per-agent quality scores"),
        ("security-reviewer", "remediation list"),
        ("git-expert", "commit → push → PR"),
    ]
    y = 1.65
    row_h = 0.42
    for i, (name, desc) in enumerate(stages):
        # connector line
        if i > 0:
            line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Inches(1.4), Inches(y - 0.05),
                                       Inches(0.04), Inches(0.07))
            line.fill.solid()
            line.fill.fore_color.rgb = ACCENT_DIM
            line.line.fill.background()
        # node
        node = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(0.75), Inches(y), Inches(3.6), Inches(row_h))
        node.fill.solid()
        node.fill.fore_color.rgb = PANEL
        node.line.color.rgb = ACCENT
        node.line.width = Pt(1)
        node.shadow.inherit = False
        tf = node.text_frame
        tf.margin_left = Inches(0.15)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = name
        r.font.name = "Consolas"
        r.font.size = Pt(13)
        r.font.bold = True
        r.font.color.rgb = ACCENT

        add_textbox(s, Inches(4.5), Inches(y + 0.08), Inches(8), Inches(row_h),
                    desc, size=13, color=TEXT)
        y += row_h + 0.05

    add_notes(s, (
        "This is the slide everything else hangs off. Don't rush.\n\n"
        "Walk it once at a high level:\n"
        "'Idea on the left, PR on the right. Everything in between is automated, "
        "but every stage produces a named artifact you can read and audit.'\n\n"
        "Then walk it again calling out the four moments that matter:\n\n"
        "  1. interrogator -> spec-drafter: 'This is where vague becomes specific. "
        "You can't build the wrong thing if you couldn't articulate it.'\n"
        "  2. orchestrator -> architect: 'Plan before code. The agent decomposes the work "
        "and proposes architecture BEFORE touching files.'\n"
        "  3. implementation-loop: 'Five iterations max. Tests drive the loop. "
        "If they don't pass, the loop runs again with the failure as context.'\n"
        "  4. karen + security-reviewer: 'Two independent audits at the end. "
        "karen checks did we build what was asked. security checks did we ship anything dangerous. "
        "Either can block the PR.'\n\n"
        "Inevitable question: 'Who's Karen?' — yes, she's named for the meme. "
        "She audits ruthlessly and returns PASS, PARTIAL, or FAIL.\n\n"
        "Optional aside if you have time: every arrow in this diagram corresponds to "
        "a markdown file in .claude/agents/. The whole pipeline is data — "
        "no hardcoded orchestration logic. That's why it's portable."
    ))


def slide_agents(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, prs)
    add_slide_header(s, "12 agents, one job each")
    rows = [
        ["Agent", "Role"],
        ["orchestrator", "Decompose spec into tasks"],
        ["coder", "Implement one task"],
        ["karen", "Audit work against spec"],
        ["architect", "Draft / update architecture docs"],
        ["spec-drafter", "Turn interrogation into SPEC.md"],
        ["concern-resolver", "Surface security concerns from triggers"],
        ["security-reviewer", "Final security pass on the diff"],
        ["git-expert", "Branch, commit, push, PR"],
        ["test-runner", "Run suite, report pass/fail/blocked"],
        ["unit-test-writer", "Generate tests to ≥90% coverage"],
        ["session-keeper", "Maintain PROGRESS_TRACKER.md"],
        ["agent-evaluator", "Score agents against criteria.json"],
    ]
    add_table(s, Inches(0.75), Inches(1.65), Inches(11.8), Inches(4.7), rows,
              col_widths=[1, 2.5], font_size=12)

    add_panel(s, Inches(0.75), Inches(6.0), Inches(11.8), Inches(0.7))
    add_rich_paragraphs(s, Inches(1.0), Inches(6.1), Inches(11.5), Inches(0.5), [
        [("Rule:", {"bold": True, "color": ACCENT, "size": 16}),
         (" if an agent's description needs the word ", {"size": 16}),
         ("\"and,\"", {"italic": True, "color": WARN, "size": 16}),
         (" split it.", {"size": 16})],
    ])
    add_notes(s, (
        "Don't read the table. Let the audience scan it.\n\n"
        "What to say while they scan:\n"
        "'Twelve agents. Every one of them is under 200 lines of prompt. "
        "Every one has a single responsibility. If you've ever written a CLAUDE.md "
        "that's two thousand lines long trying to teach one giant agent every skill — "
        "this is the alternative. Small, composable, replaceable.'\n\n"
        "Then land the rule at the bottom hard:\n"
        "'If you can't describe what an agent does without using the word and, "
        "you have two agents pretending to be one. Split them.'\n\n"
        "Call out karen and agent-evaluator specifically as the meta-agents — "
        "they review the work of other agents. That's how the system catches itself "
        "when individual agents go off the rails.\n\n"
        "If asked which agent was hardest to write: orchestrator. "
        "Task decomposition is the place where vagueness becomes expensive."
    ))


def slide_hooks(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, prs)
    add_slide_header(s, "Always-on safety rails (hooks)")
    rows = [
        ["Hook", "Fires on", "What it guarantees"],
        ["log_tool_call.sh", "Every tool call", "JSONL audit trail of every action"],
        ["secrets-postwrite.sh", "Write / Edit", "Scans for leaked credentials"],
        ["test-immutability.sh", "Edit / Write", "BLOCKS edits to existing tests"],
        ["git-commit-scope.sh", "git commit", "Injects git diff --stat first"],
        ["session-autocommit.sh", "Tracker writes", "Auto-commits progress (team mode)"],
    ]
    add_table(s, Inches(0.5), Inches(1.65), Inches(12.3), Inches(3.5), rows,
              col_widths=[1.2, 1, 2], font_size=13)

    add_panel(s, Inches(0.75), Inches(5.5), Inches(11.8), Inches(1.3))
    add_rich_paragraphs(s, Inches(1.0), Inches(5.7), Inches(11.5), Inches(1.1), [
        [("These fire whether the agent remembers them or not.", {"bold": True, "color": ACCENT, "size": 20})],
        [("", {"size": 6})],
        [("Hooks live in .claude/settings.json — the runtime enforces them, not the agent.",
          {"size": 14, "color": MUTED})],
    ])
    add_notes(s, (
        "This is the slide where security and SRE people lean in. Sell it hard.\n\n"
        "For each hook, one sentence:\n\n"
        "  - log_tool_call.sh: 'Every Bash command, every file edit, every web fetch — "
        "appended as JSONL. When the agent does something weird, you have a transcript.'\n"
        "  - secrets-postwrite.sh: 'The moment a file is written, it's scanned. "
        "If a credential lands, you know before you push.'\n"
        "  - test-immutability.sh: 'This is the one I'm proudest of. Agents are VERY good "
        "at making failing tests pass by editing the test. This hook blocks the edit at the "
        "tool-call level. New tests are fine. Existing tests are immutable.'\n"
        "  - git-commit-scope.sh: 'Diff before commit. Every time. The agent sees what "
        "it's about to commit and can bail if the scope is wrong.'\n"
        "  - session-autocommit.sh: 'Optional team mode — the agent's progress tracker "
        "auto-commits, so reviewers can watch the run live.'\n\n"
        "The closer is the key sentence: 'These fire whether the agent remembers them or not.' "
        "Hooks live in .claude/settings.json — they're harness-level guarantees, "
        "not agent-level intentions. The agent literally cannot bypass them; "
        "the runtime executes them."
    ))


def slide_gates(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, prs)
    add_slide_header(s, "Quality gates that block bad PRs")

    gates = [
        ("karen", "Audits the final diff against SPEC.md",
         "PASS → continue.  PARTIAL / FAIL → punch list to PROBLEMS.md, loop again.",
         ACCENT),
        ("security-reviewer", "Merges app-type threat models with code findings",
         "Loads checklists from security/profiles/ based on declared app_types (web_app, api, rag, ai_agent, database, frontend, …).",
         WARN),
        ("evaluate-run", "Scores each agent against its criteria.json",
         "Writes per-agent verdicts to sessions/<run_id>/traces/<agent>.json.",
         DANGER),
    ]
    y = 1.65
    for name, sub, body, color in gates:
        add_panel(s, Inches(0.75), Inches(y), Inches(11.8), Inches(1.45))
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(y),
                                   Inches(0.1), Inches(1.45))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        add_textbox(s, Inches(1.0), Inches(y + 0.1), Inches(11.4), Inches(0.5),
                    name, size=20, bold=True, color=color, mono=True)
        add_textbox(s, Inches(1.0), Inches(y + 0.55), Inches(11.4), Inches(0.4),
                    sub, size=14, color=TEXT)
        add_textbox(s, Inches(1.0), Inches(y + 0.95), Inches(11.4), Inches(0.5),
                    body, size=12, color=MUTED)
        y += 1.55

    add_textbox(s, Inches(0.75), Inches(6.5), Inches(11.8), Inches(0.5),
                "No PR opens until all three are satisfied.",
                size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_notes(s, (
        "This is the 'trust' slide. You're answering the question: "
        "'How do I know the agent didn't ship garbage?'\n\n"
        "Three gates, three answers:\n\n"
        "  - karen: 'Did we build what the spec asked for? Not does it run — does it match. "
        "She compares the diff to SPEC.md and returns a verdict. If she says PARTIAL or FAIL, "
        "the implementation loop runs again with her punch list as context.'\n"
        "  - security-reviewer: 'Did we ship anything dangerous? She's loaded with threat models "
        "specific to the app type. If you declared this is a web_app and an api, she'll check "
        "for OWASP top 10 plus API-specific risks. The checklists merge — no manual selection.'\n"
        "  - evaluate-run: 'Did each agent actually do its job well? This is the meta-gate — "
        "agent-evaluator scores every agent's output against its own criteria. "
        "It's how the harness detects when an individual agent is regressing.'\n\n"
        "Close with the line on the slide. No PR opens until all three are happy. That's the contract."
    ))


def slide_session(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, prs)
    add_slide_header(s, "Anatomy of a session", "Everything for one run, in one folder")

    add_panel(s, Inches(0.75), Inches(1.7), Inches(11.8), Inches(3.9))
    tree = [
        "sessions/20260515-1430/",
        "  ├── checkpoint.json         current pipeline stage + metadata",
        "  ├── SPEC.md                 what we're building",
        "  ├── PLAN.md                 how, broken into tasks",
        "  ├── ARCHITECTURE.md         proposed design",
        "  ├── SECURITY_CONCERNS.md    triggered concerns + checklists",
        "  ├── PROGRESS_TRACKER.md     per-agent I/O log (auto-committed)",
        "  ├── PROBLEMS.md             karen + security findings",
        "  ├── traces/<agent>.json     per-agent evaluation verdicts",
        "  └── session_log.json        structured tool-call log",
    ]
    add_textbox(s, Inches(1.0), Inches(1.85), Inches(11.3), Inches(3.7),
                "\n".join(tree), size=14, color=TEXT, mono=True)

    # bottom row of "ops"
    ops = [
        ("Cleanup", "rm -rf sessions/<id>", DANGER),
        ("Debugging", "ls sessions/<id>", WARN),
        ("Resume", "/resume <id>", ACCENT),
    ]
    x = 0.75
    w = 3.93
    for label, cmd, color in ops:
        add_panel(s, Inches(x), Inches(5.85), Inches(w), Inches(1.0))
        add_textbox(s, Inches(x + 0.15), Inches(5.95), Inches(w - 0.3), Inches(0.4),
                    label.upper(), size=11, bold=True, color=MUTED)
        add_textbox(s, Inches(x + 0.15), Inches(6.3), Inches(w - 0.3), Inches(0.5),
                    cmd, size=16, bold=True, color=color, mono=True)
        x += w + 0.05
    add_notes(s, (
        "This slide answers: 'Where does all this stuff live?'\n\n"
        "The point to land: everything for one run is in one folder. "
        "Not scattered across docs/, reports/, .claude/state/, and the git log. One folder.\n\n"
        "That single-root design is WHY cleanup, debugging, and resume are all trivial one-liners. "
        "The bottom row of the slide is the payoff.\n\n"
        "If you have time, point at PROGRESS_TRACKER.md specifically: in team mode that file "
        "auto-commits on every write, so a reviewer can pull the branch and watch the agent's "
        "progress as if they were pair programming.\n\n"
        "Optional: mention the two session modes (local-only vs collaborative) briefly. "
        "Local-only adds three lines to .git/info/exclude. Collaborative leaves them out "
        "and lets the auto-commit hook do its thing.\n\n"
        "Transition: 'OK — how do you actually use this on your own repo?'"
    ))


def slide_adoption(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, prs)
    add_slide_header(s, "Adopting it", "Five steps. Two are optional.")
    steps = [
        ("1", "Copy .claude/, sessions/, tests/, scripts/, .github/, .gitignore into your repo", False),
        ("2", "Fill in CLAUDE.md — project name, stack, test command", False),
        ("3", "(Optional) Declare app_types in CLAUDE.md to load threat models", True),
        ("4", "(Optional) Add sessions/*/ to .git/info/exclude for local-only mode", True),
        ("5", "/run and describe what you want to build", False),
    ]
    y = 1.7
    for num, text, optional in steps:
        color = MUTED if optional else ACCENT
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.75), Inches(y),
                                      Inches(0.6), Inches(0.6))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        circle.shadow.inherit = False
        tf = circle.text_frame
        tf.margin_left = Inches(0)
        tf.margin_right = Inches(0)
        tf.margin_top = Inches(0)
        tf.margin_bottom = Inches(0)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = num
        r.font.size = Pt(20)
        r.font.bold = True
        r.font.color.rgb = BG
        add_textbox(s, Inches(1.55), Inches(y + 0.1), Inches(11), Inches(0.5),
                    text, size=16, color=TEXT if not optional else MUTED)
        y += 0.85

    add_panel(s, Inches(0.75), Inches(6.1), Inches(11.8), Inches(0.85))
    add_rich_paragraphs(s, Inches(1.0), Inches(6.25), Inches(11.5), Inches(0.6), [
        [("That's it. ", {"bold": True, "color": ACCENT, "size": 18}),
         ("No runtime to install. No service to deploy. Pure scaffold.", {"size": 18})],
    ])
    add_notes(s, (
        "Make this feel small. The barrier to adoption is the message.\n\n"
        "'Five steps. Two of them are optional. The other three are copy files, "
        "fill in a markdown file, and run a command. If you've used Claude Code at all, "
        "you can adopt this in fifteen minutes.'\n\n"
        "Specifically call out the no-runtime line. People hear 'harness' and brace for "
        "a Docker image, a server, a config file with seventeen YAML keys. "
        "None of that. It's markdown and bash.\n\n"
        "If anyone asks about CI: the .github/ folder ships with workflows for the standard "
        "hooks (secrets scan, test immutability) running on PRs too, so the guarantees "
        "extend beyond the local agent loop. Mention this only if asked — "
        "it's not the main beat."
    ))


def slide_comparison(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, prs)
    add_slide_header(s, "Plain Claude Code vs EFF-IT")
    rows = [
        ["Concern", "Plain Claude Code", "EFF-IT"],
        ["One agent or many?", "One, doing everything", "12 single-purpose"],
        ["Session crashes?", "Start over", "/resume from checkpoint"],
        ["Did we build the right thing?", "You check", "karen checks"],
        ["Security?", "Your job", "Profile-driven, every run"],
        ["Tests \"adjusted\" to pass?", "Possible", "Blocked at hook level"],
        ["Audit trail?", "Transcript only", "JSONL log + per-agent traces"],
        ["Scope drift?", "Possible", "Diff-before-commit hook"],
        ["Portable across repos?", "Per-repo CLAUDE.md", "Drop-in scaffold"],
    ]
    add_table(s, Inches(0.5), Inches(1.65), Inches(12.3), Inches(5.2), rows,
              col_widths=[1.4, 1.4, 1.4], font_size=13)
    add_notes(s, (
        "This is the closer before the demo. Don't walk every row — let people read.\n\n"
        "What to say while they read:\n"
        "'Plain Claude Code is incredibly capable. Nothing on this slide is a knock against it. "
        "The point is that everything in the right column is WORK — work you'd have to do "
        "yourself, every project, every session, with raw Claude Code. EFF-IT just bundles "
        "that work into a scaffold so you don't have to redo it.'\n\n"
        "If pressed on which row matters most, pick two:\n"
        "  - 'Tests adjusted to pass' — because it's invisible until production.\n"
        "  - 'Session crashes' — because it's the single biggest productivity killer "
        "in long agent runs.\n\n"
        "Transition to demo:\n"
        "'Let me show you what it looks like in practice.'"
    ))


def slide_demo(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, prs)
    add_slide_header(s, "Live demo", "Pick one")

    demos = [
        ("/run \"add a JWT refresh endpoint to the auth service\"", "Full pipeline. ~8 min if network cooperates.", ACCENT),
        ("/fast-lane \"rename UserService.findById to UserService.getById\"", "Skip the spec. Safer choice for short slots.", WARN),
        ("/resume 20260520-1145", "Rescue option. Use a known-good fixture session.", DANGER),
    ]
    y = 1.65
    for cmd, note, color in demos:
        add_panel(s, Inches(0.75), Inches(y), Inches(11.8), Inches(1.05))
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(y),
                                   Inches(0.1), Inches(1.05))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        add_textbox(s, Inches(1.0), Inches(y + 0.1), Inches(11.5), Inches(0.5),
                    cmd, size=15, bold=True, color=color, mono=True)
        add_textbox(s, Inches(1.0), Inches(y + 0.6), Inches(11.5), Inches(0.4),
                    note, size=13, color=MUTED)
        y += 1.15

    add_panel(s, Inches(0.75), Inches(5.3), Inches(11.8), Inches(1.6))
    add_textbox(s, Inches(1.0), Inches(5.4), Inches(11.5), Inches(0.4),
                "What to watch for:", size=14, bold=True, color=ACCENT)
    watch = [
        "•  The Q&A loop in idea-interrogator (it really does push back)",
        "•  karen returning PARTIAL the first iteration (she usually does)",
        "•  Hooks firing in the transcript margin",
    ]
    add_textbox(s, Inches(1.0), Inches(5.75), Inches(11.5), Inches(1.0),
                "\n".join(watch), size=14, color=TEXT)
    add_notes(s, (
        "Pre-flight checklist before going live:\n"
        "  1. Repo is on a clean branch (the harness will create a new one).\n"
        "  2. CLAUDE.md has app_types declared, or be ready to explain why concerns are keyword-only.\n"
        "  3. Have a fixture session ready under sessions/ in case the live run stalls — "
        "you can pivot to /resume on the fixture.\n"
        "  4. Terminal font is large enough to read from the back of the room.\n\n"
        "Demo choice guidance:\n"
        "  - /run is the most impressive but also the longest. Pick this if you have 8+ minutes.\n"
        "  - /fast-lane is the safer choice for a short slot. Skip the interrogation.\n"
        "  - /resume is the rescue option — use this if live runs are flaky.\n\n"
        "Narrate as the agent works:\n"
        "  - When interrogator asks a question: 'Notice it's pushing back — it won't draft a spec "
        "until the ambiguity is resolved.'\n"
        "  - When the coder starts editing: 'Watch the secrets hook fire on every write — "
        "there's the JSONL line appearing in the log.'\n"
        "  - When karen returns: 'There's the audit. If she says PARTIAL, we loop.'\n\n"
        "Have a fallback: if the demo dies mid-flight, switch to showing "
        "sessions/<a-completed-run>/ in the file tree. The artifacts tell the story even "
        "without the live run."
    ))


def slide_qa(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, prs)
    add_slide_header(s, "Q&A", "Pre-baked answers to the obvious questions")

    qa = [
        ("Why not just use one big agent with a great prompt?",
         "Single agents drift on long tasks. Decomposition + audit catches drift early."),
        ("What if I disagree with karen?",
         "Override and continue — her verdict is advisory at the harness level. You're still the human."),
        ("Does this lock me into Claude?",
         "Hooks and pipeline are model-agnostic. Agents are markdown — portable to any agent runtime."),
        ("What's the failure mode to worry about?",
         "Stale agents. If one hasn't been invoked in 30+ days, it's a candidate for removal."),
    ]
    y = 1.65
    for q, a in qa:
        add_panel(s, Inches(0.75), Inches(y), Inches(11.8), Inches(1.15))
        add_textbox(s, Inches(1.0), Inches(y + 0.1), Inches(11.5), Inches(0.5),
                    "Q  " + q, size=15, bold=True, color=ACCENT)
        add_textbox(s, Inches(1.0), Inches(y + 0.6), Inches(11.5), Inches(0.5),
                    "A  " + a, size=13, color=TEXT)
        y += 1.25

    add_textbox(s, Inches(0.75), Inches(6.8), Inches(11.8), Inches(0.4),
                "Thanks. EFF-IT — let the agents handle it.",
                size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_notes(s, (
        "Read the room. If hands go up immediately, you've nailed the talk. "
        "If silence, prime the pump with one of the pre-baked questions on the slide.\n\n"
        "Hard questions you might get:\n\n"
        "  - 'Have you measured productivity gains?' — Be honest. 'No formal A/B, but the runs "
        "I would have abandoned without /resume are the strongest signal. Anecdotally, "
        "idea-to-PR for a small feature is under an hour.'\n\n"
        "  - 'What about cost? 12 agents per feature must burn tokens.' — Yes. The tradeoff is "
        "fewer failed runs and fewer round trips with humans for spec clarification. "
        "Net cost vs raw Claude Code is roughly comparable for non-trivial features; "
        "worse for trivial ones (use /fast-lane).\n\n"
        "  - 'Could the agent disable a hook?' — In principle, by editing .claude/settings.json. "
        "In practice, secrets-postwrite would scan the edit, and the agent has no instruction "
        "to do this. But: don't run this scaffold against a repo with sensitive secrets you "
        "haven't rotated. Defense in depth, not magic.\n\n"
        "  - 'Is this open source?' / 'How do I contribute?' — Insert your answers based on "
        "the repo's actual license status.\n\n"
        "Close with a callback to the title."
    ))


# ---------- v2 builders: revised slides 1-4 ----------


def slide_title_v2(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, prs)
    add_textbox(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(1.5),
                "EFF-IT", size=96, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    # working-name chip directly below the wordmark
    add_rich_paragraphs(s, Inches(0.5), Inches(3.45), Inches(12.3), Inches(0.5), [
        [("(working name)", {"size": 16, "italic": True, "color": MUTED})],
    ]).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    # Subtitle moved down a touch to make room for the chip
    add_textbox(s, Inches(0.5), Inches(4.05), Inches(12.3), Inches(0.6),
                "A Claude Code harness for vibe-driven software delivery",
                size=22, color=TEXT, align=PP_ALIGN.CENTER)
    add_panel(s, Inches(2.5), Inches(5.0), Inches(8.3), Inches(1.0))
    add_textbox(s, Inches(2.5), Inches(5.15), Inches(8.3), Inches(0.7),
                "“You know what? EFF-IT. Let the agents handle it.”",
                size=20, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
                "lizhang  ·  2026", size=12, color=MUTED,
                align=PP_ALIGN.CENTER)
    add_notes(s, (
        "Open by saying the name out loud, then own the joke:\n"
        "'I'm introducing EFF-IT — sorry, it's a working name.'\n"
        "Pause. Let people laugh or roll eyes. The (working name) tag is "
        "right there on the slide so you don't have to belabor it.\n\n"
        "Then the substance:\n"
        "'This is a scaffold I built on top of Claude Code. It takes the "
        "vibe-coding experience — describe what you want, let the AI build it — "
        "and wraps it with the structure you'd normally have to remember to add "
        "yourself: specs, audits, security review, checkpoints, the works.'\n\n"
        "Don't dive into architecture yet. The next slide defines vibe coding "
        "before we critique it."
    ))


def slide_what_is_vibe_v2(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, prs)
    add_slide_header(s, "What is \"vibe coding\"?")

    # Definition panel
    add_panel(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.0))
    add_rich_paragraphs(s, Inches(0.8), Inches(1.6), Inches(11.7), Inches(0.8), [
        [("Many definitions, ", {"size": 17}),
         ("variable amounts of rigor.", {"size": 17, "italic": True, "color": WARN}),
         ("  The shape that's stuck:", {"size": 17})],
        [("an AI coding agent  +  your idea in plain text.  Nothing else required.",
          {"size": 16, "color": ACCENT, "bold": True})],
    ])

    # Promise callout
    add_panel(s, Inches(0.5), Inches(2.7), Inches(12.3), Inches(0.9), color=PANEL)
    add_rich_paragraphs(s, Inches(0.8), Inches(2.85), Inches(11.7), Inches(0.7), [
        [("The promise:  ", {"size": 18, "bold": True, "color": ACCENT}),
         ("ship features at the speed of thought. You don't need to know how it works.",
          {"size": 18, "color": TEXT})],
    ])

    # Honest caveat
    add_panel(s, Inches(0.5), Inches(3.8), Inches(12.3), Inches(2.6))
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(3.8),
                              Inches(0.1), Inches(2.6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = WARN
    bar.line.fill.background()
    bar.shadow.inherit = False
    add_textbox(s, Inches(0.8), Inches(3.95), Inches(11.7), Inches(0.5),
                "Great for prototyping.  Fragile beyond it.",
                size=20, bold=True, color=WARN)
    add_textbox(s, Inches(0.8), Inches(4.45), Inches(11.7), Inches(0.4),
                "Without guardrails, quality standards, tests, and rigor, raw vibe coding accrues:",
                size=14, color=MUTED)
    # pill row of named costs
    pill_labels = [
        "Technical debt",
        "Strange design choices",
        "Inconsistent naming",
        "Duplicate / dead logic",
    ]
    x = 0.8
    for label in pill_labels:
        _, used_w = add_pill(s, Inches(x), Inches(5.05), Inches(0.55),
                              label, color=WARN, fill=BG, size=13)
        x += used_w.inches + 0.2
    add_textbox(s, Inches(0.8), Inches(5.85), Inches(11.7), Inches(0.4),
                "That's the gap EFF-IT is built to close.",
                size=14, color=ACCENT)
    add_notes(s, (
        "Anchor the definition by saying it like the slide reads:\n"
        "'There are many definitions of vibe coding, with variable amounts of "
        "rigor. The shape that's stuck is the simple one: an AI coding agent "
        "plus your idea in plain text. That's the whole input.'\n\n"
        "Sell the promise without irony first:\n"
        "'The pitch is real — ship features at the speed of thought. You don't "
        "need to know how the code works to get something running.'\n\n"
        "Then pivot. Land the WARN panel hard:\n"
        "'It's great for prototyping. It's fragile beyond that. Without "
        "guardrails, you accrue technical debt, strange design choices, "
        "inconsistent naming, duplicate or dead logic. The code runs. It also "
        "rots.'\n\n"
        "Closing line is the bridge to slide 3: 'That's the gap EFF-IT is "
        "built to close.'"
    ))


def slide_problem_v2(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, prs)
    add_slide_header(s, "The problem with raw vibe coding",
                     "Four failure modes that show up in every long session")

    # 2x2 grid of themed cards
    card_w = Inches(6.05)
    card_h = Inches(2.15)
    gap = 0.2
    col1_x = Inches(0.5)
    col2_x = Inches(0.5 + 6.05 + gap)
    row1_y = Inches(1.5)
    row2_y = Inches(1.5 + 2.15 + gap)

    add_card(s, col1_x, row1_y, card_w, card_h,
             accent_color=DANGER, title="CONTEXT LOSS",
             body_lines=[
                 "Forgets what it was building mid-feature",
                 "Context window overflows + compaction kicks in",
                 "Re-explores the codebase every new session",
             ])
    add_card(s, col2_x, row1_y, card_w, card_h,
             accent_color=WARN, title="QUALITY DRIFT",
             body_lines=[
                 "Tests get “adjusted” to pass (if written at all)",
                 "No audit that the feature actually matches the ask",
                 "No planning, architecture, or best practices baked in",
             ])
    add_card(s, col1_x, row2_y, card_w, card_h,
             accent_color=WARN, title="PROCESS DRIFT",
             body_lines=[
                 "Scope creeps silently mid-build",
                 "“Claude, what do you think?” → instant hype",
                 "New idea lands in the plan with no rigor",
             ])
    add_card(s, col2_x, row2_y, card_w, card_h,
             accent_color=DANGER, title="NO MEMORY",
             body_lines=[
                 "Silent assumptions, re-made every session",
                 "No record of WHY anything was chosen",
                 "Secrets hardcoded → land in commits",
             ])

    # Bottom pull-quote (the scope-creep anecdote)
    add_panel(s, Inches(0.5), Inches(6.15), Inches(12.3), Inches(0.85))
    add_rich_paragraphs(s, Inches(0.8), Inches(6.27), Inches(11.7), Inches(0.65), [
        [("“", {"size": 22, "color": ACCENT, "bold": True}),
         ("I ask Claude what it thinks, it hypes me up like I'm the next best "
          "thing in the world, and then it gets added to the plan mid-session.",
          {"size": 14, "italic": True, "color": TEXT}),
         ("”", {"size": 22, "color": ACCENT, "bold": True})],
    ])
    add_notes(s, (
        "Four failure modes, not five problems — the grouping matters. "
        "Walk each card briefly:\n\n"
        "CONTEXT LOSS: 'Sessions die. The window fills. When you come back, "
        "Claude re-explores everything from scratch — and the assumptions it "
        "makes the second time aren't the same as the first.'\n\n"
        "QUALITY DRIFT: 'Tests get adjusted to make the agent's code pass — "
        "if tests get written at all. There's no audit gate at the end. No "
        "one's asking did we build the right thing?'\n\n"
        "PROCESS DRIFT (this is the live one — share the anecdote):\n"
        "'This is the one that happens to me most. I'm in the middle of a "
        "build, an idea pops into my head, I ask Claude what it thinks. "
        "It hypes me up — yes, brilliant, let's add it — and now there's "
        "a whole new feature in the plan with zero design thought. "
        "Just because it sounded cool to do.'\n\n"
        "Read the pull-quote at the bottom — it lands harder when you say "
        "it the same way it reads.\n\n"
        "NO MEMORY: 'Silent assumptions every session. No record of WHY any "
        "choice was made. And the classic — credentials hardcoded as a "
        "convenience, ending up in a public commit.'\n\n"
        "Bridge to slide 4: 'That's the challenge EFF-IT is built to resolve.'"
    ))


def slide_what_is_eff_it_v2(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, prs)
    add_slide_header(s, "What EFF-IT is",
                     "A scaffold wrapped around your vibe-coding session")

    # Subtitle line under the header explaining the shape
    add_textbox(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.5),
                "Turns it into a workflow that is:",
                size=15, color=MUTED, align=PP_ALIGN.CENTER)

    # Five pill chips — the five adjectives
    pill_labels = ["Deterministic", "Rigorous", "Spec-driven",
                   "Test-driven", "Observable"]
    # measure total width first to center the row
    char_w = 18 * 0.075
    total_w = sum(char_w * len(lbl) + 0.5 for lbl in pill_labels) + 0.25 * (len(pill_labels) - 1)
    x = (13.333 - total_w) / 2
    for label in pill_labels:
        _, used_w = add_pill(s, Inches(x), Inches(2.1), Inches(0.65),
                              label, color=ACCENT, fill=PANEL, size=18,
                              padding=0.25)
        x += used_w.inches + 0.25

    # Stack banner
    add_panel(s, Inches(0.5), Inches(3.2), Inches(12.3), Inches(1.4))
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(3.2),
                              Inches(0.1), Inches(1.4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    bar.shadow.inherit = False
    add_textbox(s, Inches(0.8), Inches(3.35), Inches(11.7), Inches(0.5),
                "Bash  ·  Markdown  ·  GitHub Actions",
                size=22, bold=True, color=ACCENT, mono=True)
    add_textbox(s, Inches(0.8), Inches(3.9), Inches(11.7), Inches(0.6),
                "No runtime to install. Stack-agnostic — Python, TypeScript, Go, anything. "
                "It targets the host project's stack without interfering.",
                size=14, color=TEXT)

    # Compact "where it lives" reference
    add_panel(s, Inches(0.5), Inches(4.85), Inches(12.3), Inches(1.85))
    add_textbox(s, Inches(0.8), Inches(4.95), Inches(11.7), Inches(0.4),
                "WHERE IT LIVES", size=11, bold=True, color=MUTED)
    lives = [
        ("  .claude/commands/", "slash workflows  —  /run, /fast-lane, /resume"),
        ("  .claude/agents/", "12 single-purpose subagents"),
        ("  .claude/hooks/", "lifecycle scripts  —  logging, secrets, test-immutability"),
        ("  sessions/<run_id>/", "self-contained per-run artifacts + checkpoint"),
        ("  security/profiles/", "app-type threat models loaded per run"),
    ]
    y = 5.35
    for path, desc in lives:
        add_textbox(s, Inches(0.8), Inches(y), Inches(4.0), Inches(0.3),
                    path, size=12, bold=True, color=ACCENT, mono=True)
        add_textbox(s, Inches(4.9), Inches(y), Inches(7.8), Inches(0.3),
                    desc, size=12, color=TEXT)
        y += 0.27
    add_notes(s, (
        "Open by stating the shape, not the structure:\n"
        "'EFF-IT puts a scaffold on top of the vibe-coding session. The "
        "session itself stays — you still describe what you want, the agent "
        "still builds it. What changes is the workflow it runs inside.'\n\n"
        "Walk the five pill chips. Don't dwell on each — they're a list:\n"
        "'Deterministic, rigorous, spec-driven, test-driven, observable. "
        "Those five words are what's different about the runs you'll see "
        "later in this deck.'\n\n"
        "Now the implementation reality:\n"
        "'The scaffold itself is just Bash, Markdown, and GitHub Actions. "
        "There's no runtime to install. Whether you're using Python, "
        "TypeScript, Go, or anything else for the project itself, EFF-IT "
        "doesn't interfere — it targets whatever stack your project already uses.'\n\n"
        "The directory list at the bottom is reference, not narration. "
        "Point at it once: 'Five directories. We'll come back to each.' "
        "Move on."
    ))


# ---------- comparison: originals (slides 1-4 before revision) ----------


def slide_originals_divider(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, prs)
    add_textbox(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(0.5),
                "REFERENCE", size=18, bold=True, color=WARN,
                align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.5), Inches(2.55), Inches(12.3), Inches(1.2),
                "Slides 1–4, original versions",
                size=40, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.5), Inches(3.8), Inches(12.3), Inches(0.5),
                "Kept for comparison against the revised versions at the top of the deck.",
                size=16, color=MUTED, align=PP_ALIGN.CENTER)

    # Map of what changed
    add_panel(s, Inches(1.5), Inches(4.6), Inches(10.3), Inches(2.3))
    add_textbox(s, Inches(1.8), Inches(4.75), Inches(9.7), Inches(0.4),
                "WHAT CHANGED", size=12, bold=True, color=ACCENT)
    changes = [
        ("Slide 1", "Added (working name) tag under the wordmark"),
        ("Slide 2", "Honest definition + named cost pills (tech debt, naming, dead logic…)"),
        ("Slide 3", "Restructured into 4 themed cards + scope-creep pull-quote"),
        ("Slide 4", "Lead with 5 adjectives, not directory structure"),
    ]
    y = 5.2
    for label, desc in changes:
        add_textbox(s, Inches(1.8), Inches(y), Inches(1.4), Inches(0.35),
                    label, size=13, bold=True, color=WARN)
        add_textbox(s, Inches(3.3), Inches(y), Inches(8.2), Inches(0.35),
                    desc, size=13, color=TEXT)
        y += 0.4
    add_notes(s, (
        "Skip this slide in live delivery — it's reference only.\n\n"
        "If anyone asks why something changed between drafts, this is the "
        "summary. The next four slides are the originals verbatim from the "
        "first draft of the deck, kept for direct comparison.\n\n"
        "If presenting from this deck for review (rather than to an "
        "audience), this is the divider to land on and pivot from."
    ))


# ---------- appendix: case study (webhooksig run) ----------


def add_picture_fit(slide, image_path, left, top, max_width, max_height,
                    border=True):
    """Insert an image preserving aspect ratio, fitted inside the bounding box
    and centered within it. Returns the picture shape."""
    with Image.open(image_path) as img:
        iw, ih = img.size
    box_w = float(max_width)
    box_h = float(max_height)
    img_ar = iw / ih
    box_ar = box_w / box_h
    if img_ar >= box_ar:
        new_w = box_w
        new_h = box_w / img_ar
    else:
        new_h = box_h
        new_w = box_h * img_ar
    cx = left + (box_w - new_w) / 2
    cy = top + (box_h - new_h) / 2
    pic = slide.shapes.add_picture(str(image_path), cx, cy,
                                    width=int(new_w), height=int(new_h))
    if border:
        pic.line.color.rgb = ACCENT_DIM
        pic.line.width = Pt(0.75)
    return pic


def add_case_study_header(slide, idx_in_appendix, title, subtitle=None):
    add_accent_bar(slide, Inches(0.5), Inches(0.45), height=Inches(0.55))
    add_textbox(slide, Inches(0.75), Inches(0.32), Inches(2.0), Inches(0.4),
                f"APPENDIX · {idx_in_appendix}", size=11, bold=True, color=ACCENT)
    add_textbox(slide, Inches(0.75), Inches(0.65), Inches(12), Inches(0.5),
                title, size=24, bold=True, color=TEXT)
    if subtitle:
        add_textbox(slide, Inches(0.75), Inches(1.05), Inches(12), Inches(0.35),
                    subtitle, size=12, color=MUTED)


def add_caption(slide, text, *, color=TEXT):
    add_panel(slide, Inches(0.5), Inches(6.35), Inches(12.3), Inches(0.55))
    add_textbox(slide, Inches(0.75), Inches(6.45), Inches(11.8), Inches(0.4),
                text, size=13, color=color, align=PP_ALIGN.LEFT)


def slide_appendix_divider(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, prs)
    # Big banner
    add_textbox(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(0.5),
                "APPENDIX", size=18, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.5), Inches(2.55), Inches(12.3), Inches(1.2),
                "Case study: /run builds webhooksig",
                size=44, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.5), Inches(3.8), Inches(12.3), Inches(0.5),
                "Idea → merged PR in 30 m 13 s. Every stage. Real receipts.",
                size=18, color=MUTED, align=PP_ALIGN.CENTER)

    # Quick stats row
    stats = [
        ("30m 13s", "wall clock", ACCENT),
        ("12", "agents invoked", WARN),
        ("182", "tests written", ACCENT),
        ("$0.05", "approx token cost", DANGER),
    ]
    x = 0.75
    card_w = 2.95
    for big, small, color in stats:
        add_panel(s, Inches(x), Inches(4.8), Inches(card_w), Inches(1.2))
        add_textbox(s, Inches(x), Inches(4.9), Inches(card_w), Inches(0.6),
                    big, size=28, bold=True, color=color,
                    align=PP_ALIGN.CENTER)
        add_textbox(s, Inches(x), Inches(5.5), Inches(card_w), Inches(0.4),
                    small, size=11, color=MUTED, align=PP_ALIGN.CENTER)
        x += card_w + 0.1

    add_textbox(s, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.4),
                "Full screen recording available alongside this deck: "
                "pptx_refs/Screen Recording 2026-05-22 at 10.51.37 AM.mov",
                size=10, color=MUTED, align=PP_ALIGN.CENTER, font="Consolas")
    add_notes(s, (
        "Pause here. This is the pivot from 'here's the harness' to "
        "'here's a real run of the harness, soup to nuts.'\n\n"
        "What to say:\n"
        "'The rest of the talk was about the harness. The rest of the deck is "
        "what one run looks like in practice. Thirty minutes, twelve agents, "
        "one merged PR, five cents in tokens. We're going to walk every stage "
        "with the actual screenshots from that session.'\n\n"
        "If you have time and the audience is engaged, mention you can also "
        "play the full screen recording — it's in pptx_refs/ next to the deck. "
        "Don't try to embed a 30 MB .mov in the slides; just have it ready in "
        "a side window if asked."
    ))


def slide_cs_invocation(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, prs)
    add_case_study_header(s, "1 / 8",
                          "Step 0–1: Invocation + interrogation",
                          "From /run to a working pitch in two minutes")
    left_img = prep_image("10.41.53", "01a-invoke")
    right_img = prep_image("10.43.26", "01b-interrogator-pitch")
    add_picture_fit(s, left_img, Inches(0.5), Inches(1.45),
                    Inches(6.1), Inches(4.75))
    add_picture_fit(s, right_img, Inches(6.75), Inches(1.45),
                    Inches(6.1), Inches(4.75))
    add_caption(s, "Left: /run skill confirmation.   "
                   "Right: the interrogator's working pitch — webhooksig was born here.")
    add_notes(s, (
        "Two screenshots, two beats:\n\n"
        "LEFT (10:41:53): 'This is what kicking off a run looks like — type /run, "
        "Claude Code prompts you to allow the skill. That's the only friction in "
        "the whole 30 minutes.'\n\n"
        "RIGHT (10:43:26): 'A minute and a half later, the interrogator has "
        "summarized the idea back to me. Note the prompt at the bottom — Is this "
        "the right one-liner, or do you want to refine it? It will not draft a "
        "spec until you confirm or correct.'\n\n"
        "Optional aside: the model badge in the bottom-left switched from "
        "Sonnet 4.6 to Opus 4.7 between these two screenshots. The skill "
        "actually suggested the swap because Opus does deeper interrogation. "
        "Small UX detail, but it tells you the harness is opinionated about "
        "model choice per stage."
    ))


def slide_cs_spec(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, prs)
    add_case_study_header(s, "2 / 8",
                          "Step 2: SPEC.md emerges",
                          "The contract karen audits against, hours from now")
    img = prep_image("10.43.59", "02-spec")
    add_picture_fit(s, img, Inches(0.5), Inches(1.45),
                    Inches(12.3), Inches(4.75))
    add_caption(s, "Goals, non-goals, users, scope — drafted before a single source file is touched.")
    add_notes(s, (
        "Don't read the spec aloud. Let people scan it for ten seconds.\n\n"
        "What to say:\n"
        "'Every section you see here came out of the interrogator's Q&A. "
        "Goals, non-goals, intended users, scope boundaries. This file becomes "
        "the contract — karen, the auditor agent at the end of the pipeline, "
        "literally diffs the implementation against this spec to decide PASS, "
        "PARTIAL, or FAIL.'\n\n"
        "Land this point: 'In raw vibe coding, the spec lives in chat history "
        "and gets lost. Here, it's a file. You can read it, share it, and "
        "the agents can re-read it whenever they drift.'"
    ))


def slide_cs_security(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, prs)
    add_case_study_header(s, "3 / 8",
                          "Step 3: Security concerns, merged per app_type",
                          "The threat-model story from slide 9 — in production form")
    left_img = prep_image("10.48.07", "03a-security-concerns-1")
    right_img = prep_image("10.48.33", "03b-security-concerns-2")
    add_picture_fit(s, left_img, Inches(0.5), Inches(1.45),
                    Inches(6.1), Inches(4.75))
    add_picture_fit(s, right_img, Inches(6.75), Inches(1.45),
                    Inches(6.1), Inches(4.75))
    add_caption(s, "concern-resolver reads security/concerns/*.md and unions per declared app_type → SECURITY_CONCERNS.md.")
    add_notes(s, (
        "These two screenshots are SECURITY_CONCERNS.md, top half on the left, "
        "bottom half on the right. Walk through what's on screen:\n\n"
        "'The concern-resolver looked at the SPEC, looked at app_types declared "
        "in CLAUDE.md, and produced this merged checklist. SR-1 is timing-safe "
        "comparison. SR-2 is replay protection. SR-3 is no-secret-in-logs. "
        "Every one of these has a triggered-by reason and a remediation hint.'\n\n"
        "The key point: this happens BEFORE implementation. The architect agent "
        "reads this file and bakes the mitigations into the design. The "
        "security-reviewer at the end reads the SAME file to verify nothing "
        "was missed. One checklist, three agents."
    ))


def slide_cs_architecture(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, prs)
    add_case_study_header(s, "4 / 8",
                          "Step 4: ARCHITECTURE.md (plan-before-code)",
                          "Trust boundary + HMAC call path, written before any source file")
    left_img = prep_image("10.49.47", "04a-arch-system-context")
    right_img = prep_image("10.50.42", "04b-arch-data-flow")
    add_picture_fit(s, left_img, Inches(0.5), Inches(1.45),
                    Inches(6.1), Inches(4.75))
    add_picture_fit(s, right_img, Inches(6.75), Inches(1.45),
                    Inches(6.1), Inches(4.75))
    add_caption(s, "Left: system context with trust boundary.   Right: HMAC verify pseudocode — the algorithm settled before the code.")
    add_notes(s, (
        "Two halves of one architecture doc.\n\n"
        "LEFT: system context diagram. Note the trust boundary — webhooksig "
        "sits INSIDE the caller's request handler, after TLS termination, "
        "before business logic. The architect agent drew this. ASCII art, "
        "not Mermaid — because it has to render in any markdown viewer, "
        "including git diff.\n\n"
        "RIGHT: the data-flow section. Pseudocode for the call path through "
        "verify_github and verify_stripe. This is the algorithm being settled "
        "BEFORE any code is written. By the time the coder agent gets a task, "
        "all of these decisions are pre-made.\n\n"
        "Land it: 'Most agent runs go straight from spec to code. The architect "
        "step is what stops the coder from re-deciding architecture in every "
        "function it writes.'"
    ))


def slide_cs_parallel_coders(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, prs)
    add_case_study_header(s, "5 / 8",
                          "Step 5: Parallel coders + advisor escalation",
                          "Four coders concurrent. One hits an ambiguous edge — calls the advisor.")
    left_img = prep_image("10.58.45", "05a-parallel-coders")
    right_img = prep_image("10.59.11", "05b-advisor")
    add_picture_fit(s, left_img, Inches(0.5), Inches(1.45),
                    Inches(6.1), Inches(4.75))
    add_picture_fit(s, right_img, Inches(6.75), Inches(1.45),
                    Inches(6.1), Inches(4.75))
    add_caption(s, "Left: four agents running in parallel.   Right: one of them consults the advisor on a header-length bound.")
    add_notes(s, (
        "LEFT: 'Four agents running in parallel — fixing test mismatches and "
        "implementing the P2 module tier.' This is the orchestrator's "
        "decomposition paying off — independent tasks fan out to independent "
        "coders. Wall-clock for this phase: under five minutes.\n\n"
        "RIGHT: this is the moment that justifies the advisor tool. The coder "
        "hit an ambiguous spec edge — what's the right max_header_length bound? "
        "Instead of guessing, it consulted the stronger advisor model. The "
        "advisor reviewed the full transcript and confirmed the approach.\n\n"
        "Key point: 'The advisor isn't a fallback for failure. It's a "
        "second-opinion mechanism for decisions that have consequences. "
        "The coder used it correctly here — at a branch point, before committing "
        "to an interpretation.'"
    ))


def slide_cs_karen(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, prs)
    add_case_study_header(s, "6 / 8",
                          "Karen + human-in-the-loop triage",
                          "When the audit surfaces work outside scope, the human stays in control")
    img = prep_image("11.01.12", "06-karen-triage")
    add_picture_fit(s, img, Inches(0.5), Inches(1.45),
                    Inches(12.3), Inches(4.75))
    add_caption(s, "160 / 182 tests pass. 22 failures = pre-existing test bugs. Karen names the failure mode and proceeds. Human picked Option C.")
    add_notes(s, (
        "This is the slide where 'agents in the loop' becomes 'humans still in "
        "the loop when it counts.'\n\n"
        "Walk the screenshot:\n"
        "  - pytest result: 160 passing, 22 failing\n"
        "  - karen identifies: '22 failures — all pre-existing test bugs from "
        "    the unit-test-writer (wrong kwargs, wrong provider expectations, "
        "    test design mismatches)'\n"
        "  - Karen presents options. The human picked Option C — accept the "
        "    pre-existing failures, proceed.\n\n"
        "What to say:\n"
        "'Karen could have looped forever trying to fix unit-test-writer's bugs. "
        "Instead, she recognized the failure pattern, escalated with options, "
        "and let me make the call. This is the difference between an autonomous "
        "agent that wastes 20 minutes and a harness that knows when to ask.'\n\n"
        "Optional follow-up: the 22 pre-existing bugs got logged in PROBLEMS.md "
        "and PLAN.md for a future run to address. Nothing got swept under the rug."
    ))


def slide_cs_scoreboard(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, prs)
    add_case_study_header(s, "7 / 8",
                          "The final scoreboard",
                          "Every stage's result, one screen, baked in 30 m 13 s")
    img = prep_image("11.03.22", "07-scoreboard")
    add_picture_fit(s, img, Inches(0.5), Inches(1.45),
                    Inches(12.3), Inches(4.75))
    add_caption(s, "TDD red phase · parallel + sequential coders · 22 pre-existing fails accepted · 2 medium security findings · 1 fixed, 1 documented · PR open.",
                color=ACCENT)
    add_notes(s, (
        "This is the money shot. Don't talk over it for the first ten seconds — "
        "let people read.\n\n"
        "When you talk, hit the rows that matter:\n"
        "  - 'TDD red phase: 182 tests written BEFORE implementation. That's "
        "    the test-immutability hook earning its keep — agents wrote tests "
        "    first, then code to make them pass.'\n"
        "  - 'Parallel + sequential split: the orchestrator decided which work "
        "    was safe to parallelize. P1/P2 groups (types, core, providers) in "
        "    parallel. S1/S2 (cli, entrypoint) sequential because they depend "
        "    on the others.'\n"
        "  - 'Karen: PARTIAL. Why partial and not PASS? She flagged a "
        "    .gitignore hiding src/, plus a provider-unknown design note. "
        "    Both legitimate, both reviewable.'\n"
        "  - 'Security review: two medium items. SR-1 (UnicodeDecodeError in "
        "    Stripe) fixed in the same run. SR-2 (secret visible in process "
        "    arg list) documented in PLAN.md for follow-up. Neither shipped "
        "    silently.'\n\n"
        "Land the closer:\n"
        "'30 minutes, 13 seconds. Two commits, PR open. Real working library. "
        "This is what the green column of slide 12 actually looks like in "
        "wall-clock time.'"
    ))


def slide_cs_jsonl(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, prs)
    add_case_study_header(s, "8 / 8",
                          "Receipts under the hood",
                          "The JSONL audit trail from slide 8, in the wild")
    left_img = prep_image("10.44.33", "08a-jsonl-1")
    right_img = prep_image("10.46.30", "08b-jsonl-2")
    add_picture_fit(s, left_img, Inches(0.5), Inches(1.45),
                    Inches(6.1), Inches(4.75))
    add_picture_fit(s, right_img, Inches(6.75), Inches(1.45),
                    Inches(6.1), Inches(4.75))
    add_caption(s, "log_tool_call.sh in action: every Read, Bash, Edit, Agent spawn, web fetch — append-only JSONL.")
    add_notes(s, (
        "This slide exists for the skeptics. On slide 8 you claimed an audit "
        "trail. Here it is.\n\n"
        "What to say:\n"
        "'These are excerpts from sessions/20260520-2116/session_log.json. Every "
        "tool call the agents made during the run got appended to this file. "
        "Bash commands, file reads, file writes, sub-agent invocations, "
        "everything. When something goes weird in a run — and things do go "
        "weird — this is the file you grep.'\n\n"
        "If asked what you actually do with the log:\n"
        "  - Replay a run mentally without re-running it\n"
        "  - Find which agent touched a specific file\n"
        "  - Measure how often a given tool gets called\n"
        "  - Catch agents trying things they shouldn't (the hooks block most, "
        "    but the log shows the attempt)\n\n"
        "If the audience is non-technical, you can soft-pedal this slide: "
        "'For the engineers in the room — yes, every action is logged. Moving on.' "
        "and skip the deep-dive."
    ))


# ---------- assemble ----------

def build(output_path: Path) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    builders = [
        # main deck — v2 versions for slides 1-4
        slide_title_v2,
        slide_what_is_vibe_v2,
        slide_problem_v2,
        slide_what_is_eff_it_v2,
        slide_entry_points,
        slide_pipeline,
        slide_agents,
        slide_hooks,
        slide_gates,
        slide_session,
        slide_adoption,
        slide_comparison,
        slide_demo,
        slide_qa,
        # case-study appendix
        slide_appendix_divider,
        slide_cs_invocation,
        slide_cs_spec,
        slide_cs_security,
        slide_cs_architecture,
        slide_cs_parallel_coders,
        slide_cs_karen,
        slide_cs_scoreboard,
        slide_cs_jsonl,
        # originals kept for comparison
        slide_originals_divider,
        slide_title,
        slide_what_is_vibe,
        slide_problem,
        slide_what_is_eff_it,
    ]

    total = len(builders)
    for i, build_fn in enumerate(builders, start=1):
        build_fn(prs)
        slide = prs.slides[i - 1]
        if i > 1:
            add_footer(slide, i, total)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path


if __name__ == "__main__":
    here = Path(__file__).resolve().parent.parent
    default_out = here / "docs" / "EFF-IT.pptx"
    out = Path(sys.argv[1]) if len(sys.argv) > 1 else default_out
    saved = build(out)
    print(f"Wrote {saved}")
